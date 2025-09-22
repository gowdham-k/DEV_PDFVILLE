from restrictions import check_restrictions
from flask import jsonify, send_file, request
from pdf2image import convert_from_bytes
from pptx import Presentation
import traceback
import os
from utils import cleanup_file, create_temp_file, create_temp_dir
from PyPDF2 import PdfReader


def convert_pdf_to_pptx():
    """Convert a PDF to PowerPoint presentation with restrictions"""

    try:
        uploaded_file = request.files.get("file") or (
            request.files.getlist("files")[0] if request.files.getlist("files") else None
        )
        if not uploaded_file:
            return jsonify({"error": "No file provided"}), 400

        # Get email (default = free user for demo)
        email = request.form.get("email", "free@example.com")

        # Save the uploaded file temporarily
        temp_dir = create_temp_dir()
        file_path = os.path.join(temp_dir, uploaded_file.filename)
        uploaded_file.save(file_path)

        # Restriction check
        restriction = check_restrictions(email, [file_path])
        if restriction:
            return jsonify({"error": restriction}), 403

        # Convert PDF pages to images
        with open(file_path, "rb") as f:
            images = convert_from_bytes(f.read())

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        temp_img_paths = []
        for img in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            img_path = create_temp_file(".jpg")
            img.save(img_path, "JPEG")
            temp_img_paths.append(img_path)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        output_path = create_temp_file(".pptx")
        prs.save(output_path)

        # Cleanup temp images
        for img_path in temp_img_paths:
            cleanup_file(img_path)

        return send_file(output_path, as_attachment=True, download_name="converted.pptx")

    except Exception as e:
        print(f"Error in convert_pdf_to_pptx: {str(e)}")
        print(traceback.format_exc())
        return jsonify({"error": f"Failed to convert PDF to PPTX: {str(e)}"}), 500
